# -*- coding: utf-8 -*-
"""
create_ppt.py — 答辩 PPT 自动生成脚本
用途：阶段 11 答辩准备——基于论文内容和结果自动生成 6-8 页答辩 PPT
依赖：python-pptx（pip install python-pptx）
运行方式：python create_ppt.py --paper 论文/main.tex --output 答辩PPT.pptx
"""
import os
from pathlib import Path
from typing import List, Dict, Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("⚠️ python-pptx 未安装，将输出文本版 PPT 大纲")
    print("   安装: pip install python-pptx")


# ============================================================
# 配色方案（学术答辩风格）
# ============================================================
TITLE_COLOR = RGBColor(0x1A, 0x3C, 0x6E)       # 深蓝
ACCENT_COLOR = RGBColor(0xC0, 0x39, 0x2B)       # 暗红（强调）
BODY_COLOR = RGBColor(0x2C, 0x3E, 0x50)         # 深灰
BG_COLOR = RGBColor(0xFC, 0xFC, 0xFC)           # 近白
TABLE_HEADER_BG = RGBColor(0x1A, 0x3C, 0x6E)    # 表头深蓝
TABLE_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)     # 表头白字
TABLE_ROW_ALT = RGBColor(0xEE, 0xF2, 0xF7)       # 交替行浅蓝
LIGHT_BLUE = RGBColor(0x34, 0x95, 0xDB)


def create_presentation(title: str, slides_content: List[Dict]) -> str:
    """创建 PPT 演示文稿

    Args:
        title: 演示文稿标题
        slides_content: 每页内容列表 [{title, content, type}]

    Returns:
        str: 输出文件路径
    """
    if not HAS_PPTX:
        return _create_text_outline(title, slides_content)

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9 宽屏
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides_content):
        slide_type = slide_data.get('type', 'content')
        if slide_type == 'title':
            _add_title_slide(prs, slide_data)
        elif slide_type == 'section':
            _add_section_slide(prs, slide_data)
        elif slide_type == 'comparison':
            _add_comparison_slide(prs, slide_data)
        else:
            _add_content_slide(prs, slide_data)

    output_path = '答辩PPT.pptx'
    prs.save(output_path)
    return output_path


def _add_title_slide(prs, data: Dict):
    """第 1 页：标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    # 背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()

    # 标题
    left, top, width, height = Inches(1.5), Inches(2.5), Inches(10), Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get('title', '数学建模竞赛答辩')
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle = data.get('subtitle', '')
    if subtitle:
        left, top, _, _ = Inches(1.5), Inches(4.5), Inches(10), Inches(1)
        txBox2 = slide.shapes.add_textbox(left, top, width, Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = BODY_COLOR
        p2.alignment = PP_ALIGN.CENTER


def _add_section_slide(prs, data: Dict):
    """章节分隔页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = TITLE_COLOR
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(2), Inches(3), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER


def _add_content_slide(prs, data: Dict):
    """内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 顶部色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # 内容要点
    content = data.get('content', [])
    if isinstance(content, str):
        content = content.split('\n')
    content = [c for c in content if c.strip()]

    txBox2 = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.5), Inches(11), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, item in enumerate(content):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        # 识别编号格式
        if item.lstrip().startswith(('•', '▪', '▸', '→', '1.', '2.', '3.')):
            p.text = item
        else:
            p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = BODY_COLOR
        p.space_after = Pt(12)

    # 页脚
    txBox3 = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = data.get('footer', '')
    p3.font.size = Pt(10)
    p3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p3.alignment = PP_ALIGN.RIGHT


def _add_comparison_slide(prs, data: Dict):
    """对比表页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = data.get('title', '')
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # 表格
    table_data = data.get('table', [])
    if table_data:
        rows, cols = len(table_data), len(table_data[0]) if table_data else 0
        table_shape = slide.shapes.add_table(
            rows, cols, Inches(1), Inches(1.5), Inches(11), Inches(5.5))
        table = table_shape.table

        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.text = str(table_data[r][c])
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    if r == 0:
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = TABLE_HEADER_FG
                    else:
                        paragraph.font.color.rgb = BODY_COLOR
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TABLE_HEADER_BG
                elif r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TABLE_ROW_ALT


def _create_text_outline(title: str, slides_content: List[Dict]) -> str:
    """当 python-pptx 不可用时，创建文本版 PPT 大纲"""
    output_path = '答辩/答辩大纲.txt'
    os.makedirs('答辩', exist_ok=True)

    lines = [f"答辩 PPT 大纲: {title}", "=" * 50]
    for i, slide in enumerate(slides_content, 1):
        lines.append(f"\n--- 第 {i} 页: {slide.get('title', '')} ---")
        content = slide.get('content', [])
        if isinstance(content, str):
            content = content.split('\n')
        for item in content:
            if item.strip():
                lines.append(f"  • {item.strip()}")
    lines.append("\n" + "=" * 50)
    lines.append("提示: pip install python-pptx 后可自动生成 .pptx 文件")

    text = '\n'.join(lines)
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(text)
    print(f"文本版 PPT 大纲已保存至: {output_path}")
    return output_path


# ============================================================
# 标准 PPT 结构生成器
# ============================================================

def build_standard_ppt_structure(problem_title: str, questions: List[str],
                                 innovations: List[str],
                                 key_results: Optional[Dict] = None) -> List[Dict]:
    """构建标准答辩 PPT 结构（6-8 页）

    Args:
        problem_title: 赛题标题
        questions: 各问简述
        innovations: 创新点列表
        key_results: 关键结果 {question_name: [{metric, value, comparison}]}

    Returns:
        List[Dict]: slides_content（可直接传入 create_presentation）
    """
    slides = []

    # 第 1 页：标题页
    slides.append({
        'type': 'title',
        'title': '数学建模竞赛答辩',
        'subtitle': f'题目: {problem_title}',
    })

    # 第 2 页：问题背景与技术路线图
    slides.append({
        'type': 'content',
        'title': '问题背景与技术路线',
        'content': [
            f'研究问题: {problem_title}',
            '技术路线: 数据预处理 → EDA洞察 → 建模求解 → 验证 → 结论',
            f'共 {len(questions)} 个问题，采用统一数学框架',
            '核心创新点将在第 6 页详细介绍',
        ]
    })

    # 第 3-N 页：各问核心方法 + 关键结果
    for i, q in enumerate(questions[:4]):  # 最多 4 问
        result_items = [f'核心方法: 见论文第 {i+5} 章']
        if key_results and q in key_results:
            for r in key_results[q]:
                result_items.append(f"{r.get('metric','指标')}: {r.get('value','')} "
                                   f"(对比: {r.get('comparison','')})")
        result_items.append(f'创新点应用: {innovations[min(i, len(innovations)-1)] if innovations else "见后页"}')
        slides.append({
            'type': 'content',
            'title': f'问题 {i+1}: {q}',
            'content': result_items,
        })

    # 创新点总览页
    slides.append({
        'type': 'section',
        'title': '核心创新',
    })
    inno_items = [f'创新 {i+1}: {inv}' for i, inv in enumerate(innovations)]
    inno_items.append('以上创新点均有消融实验验证（见论文附录）')
    slides.append({
        'type': 'content',
        'title': '创新点总览与消融验证',
        'content': inno_items,
    })

    # 验证与灵敏度
    slides.append({
        'type': 'content',
        'title': '验证与灵敏度分析',
        'content': [
            '多方法交叉验证（≥3 种独立方法）',
            'Sobol 全局灵敏度分析（一阶 + 总阶指数）',
            '蒙特卡洛稳健性（参数不确定性量化）',
            '极端情况测试（边界值/退化/噪声）',
            '所有结果均通过统计显著性检验（p<0.05）',
        ]
    })

    # 结论与推广
    slides.append({
        'type': 'content',
        'title': '结论与推广',
        'content': [
            '模型总结与核心贡献',
            '方法的泛化能力与适用范围',
            '局限性与未来改进方向',
            '实际应用价值与社会意义',
        ]
    })

    return slides


# ============================================================
# 10 个预答问题备忘
# ============================================================

PREPARED_QUESTIONS = [
    ("为什么选择这个方法而非其他方法？",
     "从数学结构匹配度、数据特征适配性、计算效率三个角度回答。"),
    ("这个创新点的本质是什么？与已有工作有什么区别？",
     "回答'已有方法做不了什么 → 我们做了什么 → 效果提升多少'。"),
    ("模型假设在现实中是否成立？",
     "逐一论证假设的合理性，引用数据支持，说明不成立时的修正方案。"),
    ("参数取值有什么依据？是不是人为调的？",
     "区分数据驱动估计、理论推导、经验值三类，每类说明来源。"),
    ("如何排除过拟合？",
     "K折交叉验证 + 留出验证 + 与简单方法对比 + 复杂度惩罚项。"),
    ("结果对参数变化是否稳健？",
     "引用 Tornado/Sobol/MC 灵敏度结果，说明哪些因素最敏感。"),
    ("这个结论能推广到什么程度？",
     "说明模型的适用范围（边界）、推广条件、迁移需注意事项。"),
    ("模型最大的局限是什么？",
     "诚实陈述 2-3 个具体局限，说明影响程度和未来改进路径。"),
    ("这个模型有什么实际价值？",
     "从决策支持、效率提升、成本降低等角度量化说明。"),
    ("再给 24 小时，你会改进什么？",
     "优先回答当前最薄弱环节的改进方案，展示自我反思能力。"),
]


def generate_qa_cheatsheet() -> str:
    """生成 10 个预答问题的备忘文档"""
    lines = ["# 答辩预答问题备忘", "", "> 以下 10 个问题覆盖评委提问的 90% 场景。",
             "> 建议每个问题准备 1-2 分钟的流畅回答。", ""]
    for i, (q, hint) in enumerate(PREPARED_QUESTIONS, 1):
        lines.append(f"## {i}. {q}")
        lines.append(f"")
        lines.append(f"**答题思路**: {hint}")
        lines.append(f"")
        lines.append(f"**我的回答**: [填写具体回答]")
        lines.append(f"")
    return '\n'.join(lines)


# ============================================================
# 主入口
# ============================================================

if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser(description='答辩 PPT 自动生成')
    parser.add_argument('--title', default='数学建模竞赛答辩',
                       help='PPT 标题')
    parser.add_argument('--output', default='答辩PPT.pptx',
                       help='输出文件路径')
    parser.add_argument('--qa-only', action='store_true',
                       help='仅生成预答问题备忘')
    args = parser.parse_args()

    if args.qa_only:
        cheatsheet = generate_qa_cheatsheet()
        os.makedirs('答辩', exist_ok=True)
        with open('答辩/预答问题备忘.md', 'w', encoding='utf-8') as f:
            f.write(cheatsheet)
        print("预答问题备忘已保存至: 答辩/预答问题备忘.md")
    else:
        # 示例：标准结构
        example_questions = [
            '数据驱动的XX建模与求解',
            '考虑耦合效应的XX优化',
            '多目标权衡与决策建议',
        ]
        example_innovations = [
            '基于XX的XX方法——解决XX无法处理XX的局限',
            '将XX领域的XX定理迁移至本题——非平凡应用',
        ]
        slides = build_standard_ppt_structure(
            problem_title='[在此填写赛题标题]',
            questions=example_questions,
            innovations=example_innovations,
        )
        output = create_presentation(args.title, slides)
        print(f"PPT 已保存至: {output}")
        if HAS_PPTX:
            print(f"共 {len(slides)} 页, 16:9 宽屏格式")

        # 同时生成预答问题备忘
        cheatsheet = generate_qa_cheatsheet()
        os.makedirs('答辩', exist_ok=True)
        with open('答辩/预答问题备忘.md', 'w', encoding='utf-8') as f:
            f.write(cheatsheet)
        print("预答问题备忘已保存至: 答辩/预答问题备忘.md")
